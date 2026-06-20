"""
PowerPoint Snapshot Export.

Generates a 2-slide leadership-ready snapshot for the currently-viewed
SKU/category: the recommended split, and the dollar tradeoff vs. the
cheapest single-region option. Same visual language as the BMPL pitch
deck (white / charcoal / burnt-orange, Times New Roman, no bullet lists,
no accent-stripe clutter) for brand consistency across the initiative.

This is a snapshot, not a live document -- it freezes the numbers at
generation time.
"""
import io
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x33, 0x37, 0x3D)
SLATE = RGBColor(0x5B, 0x61, 0x68)
SLATE_LIGHT = RGBColor(0x8A, 0x90, 0x99)
ORANGE = RGBColor(0xD2, 0x69, 0x1E)
ORANGE_DARK = RGBColor(0xB8, 0x5A, 0x18)
PANEL_GREY = RGBColor(0xF4, 0xF4, 0xF5)
BORDER_GREY = RGBColor(0xD8, 0xDA, 0xDD)
GREY_HEADER = RGBColor(0x6B, 0x71, 0x78)

FONT = "Times New Roman"


def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_shadow(shape):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    effectLst = spPr.makeelement(qn('a:effectLst'), {})
    outerShdw = spPr.makeelement(qn('a:outerShdw'), {
        'blurRad': '63500', 'dist': '25400', 'dir': '5400000', 'rotWithShape': '0'
    })
    clr = spPr.makeelement(qn('a:srgbClr'), {'val': '000000'})
    alpha = clr.makeelement(qn('a:alpha'), {'val': '12000'})
    clr.append(alpha)
    outerShdw.append(clr)
    effectLst.append(outerShdw)
    spPr.append(effectLst)


def _textbox(slide, x, y, w, h, text, size=14, color=CHARCOAL, bold=False, italic=False,
             align=PP_ALIGN.LEFT, font=FONT, anchor=None, char_spacing=None, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return tb


def _rect(slide, x, y, w, h, fill_color, line_color=None, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    _set_fill(shp, fill_color)
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
        shp.line.fill.solid()
        shp.line.fill.fore_color.rgb = line_color
    if shadow:
        _add_shadow(shp)
    return shp


def _footer(slide, label, page_num, W, H):
    _textbox(slide, 0.6, H - 0.45, 7, 0.3, label.upper(), size=9, color=SLATE_LIGHT, char_spacing=1.5)
    _textbox(slide, W - 4.9, H - 0.45, 3.7, 0.3, "VIRVENTURES   |   SHIPMENT INTELLIGENCE",
             size=9, color=SLATE_LIGHT, align=PP_ALIGN.RIGHT)
    _textbox(slide, W - 0.85, H - 0.45, 0.4, 0.3, str(page_num), size=9, color=SLATE_LIGHT, align=PP_ALIGN.RIGHT)


def _slide_title(slide, title, eyebrow):
    _textbox(slide, 0.6, 0.4, 11, 0.3, eyebrow.upper(), size=12, color=ORANGE, bold=True, char_spacing=2)
    _textbox(slide, 0.6, 0.68, 12.1, 0.7, title, size=28, color=CHARCOAL, bold=True)


def build_snapshot(sku_label, region_split: dict, demand_pct: dict, recommended_cost,
                    cheapest_region, cheapest_cost, cost_delta, coverage_gap_pct, rationale,
                    window_days, as_of_date):
    """
    Builds a 2-slide pptx in memory. Returns bytes.

    region_split: {"East": units, "Central": units, "West": units} -- recommended
    demand_pct: {"East": pct, ...} -- corrected demand basis
    coverage_gap_pct: percentage points of demand the CHEAPEST single-region
        option would miss (i.e. what the recommended split avoids missing) --
        this is recommend_split()'s coverage_gap_pct * 100, not a constant.
    """
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.333 in
    prs.slide_height = Emu(6858000)   # 7.5 in
    W, H = 13.333, 7.5
    blank_layout = prs.slide_layouts[6]

    total_units = sum(region_split.values())

    # ============== SLIDE 1: Recommended Split ==============
    s1 = prs.slides.add_slide(blank_layout)
    bg = _rect(s1, 0, 0, W, H, WHITE)
    _slide_title(s1, f"Recommended Regional Split — {sku_label}", "Shipment Intelligence Snapshot")

    _textbox(s1, 0.6, 1.35, 11, 0.35,
             f"Based on trailing {window_days}-day sell-through-corrected demand, as of {as_of_date}",
             size=12, color=SLATE, italic=True)

    card_y, card_h, card_w, gap = 2.0, 3.5, 3.7, 0.45
    regions = ["East", "Central", "West"]
    start_x = (W - (card_w * 3 + gap * 2)) / 2

    for i, region in enumerate(regions):
        x = start_x + i * (card_w + gap)
        units = region_split.get(region, 0)
        pct = demand_pct.get(region, 0) * 100
        is_dominant = units == max(region_split.values())

        _rect(s1, x, card_y, card_w, card_h, ORANGE if is_dominant else WHITE,
              line_color=None if is_dominant else BORDER_GREY, shadow=True)

        _textbox(s1, x + 0.3, card_y + 0.35, card_w - 0.6, 0.35, region.upper(),
                 size=14, color=WHITE if is_dominant else SLATE, bold=True, char_spacing=1.5,
                 align=PP_ALIGN.CENTER)

        _textbox(s1, x + 0.3, card_y + 0.95, card_w - 0.6, 1.0, f"{units:,}",
                 size=44, color=WHITE if is_dominant else CHARCOAL, bold=True, align=PP_ALIGN.CENTER)

        _textbox(s1, x + 0.3, card_y + 1.95, card_w - 0.6, 0.35, "UNITS RECOMMENDED",
                 size=10, color=WHITE if is_dominant else SLATE_LIGHT, char_spacing=1, align=PP_ALIGN.CENTER)

        _rect(s1, x + 0.5, card_y + 2.5, card_w - 1.0, 0.01,
              WHITE if is_dominant else BORDER_GREY)

        _textbox(s1, x + 0.3, card_y + 2.7, card_w - 0.6, 0.5, f"{pct:.0f}% of demand",
                 size=16, color=WHITE if is_dominant else ORANGE_DARK, bold=True, align=PP_ALIGN.CENTER)

    _textbox(s1, 0.6, card_y + card_h + 0.3, 12.1, 0.4,
             f"TOTAL SHIPMENT: {total_units:,} UNITS",
             size=13, color=CHARCOAL, bold=True, char_spacing=1, align=PP_ALIGN.CENTER)

    _footer(s1, "Recommendation", "1", W, H)

    # ============== SLIDE 2: Dollar Tradeoff ==============
    s2 = prs.slides.add_slide(blank_layout)
    _rect(s2, 0, 0, W, H, WHITE)
    _slide_title(s2, "The Cost Tradeoff, Made Explicit", "Decision Transparency")

    panel_y, panel_h = 1.9, 4.4
    left_x, left_w = 0.6, 5.9
    right_x, right_w = left_x + left_w + 0.4, 5.9

    # Left panel: demand-optimal
    _rect(s2, left_x, panel_y, left_w, panel_h, "FBEEE3" if False else RGBColor(0xFB, 0xEE, 0xE3),
          line_color=ORANGE, shadow=True)
    _textbox(s2, left_x + 0.4, panel_y + 0.35, left_w - 0.8, 0.3, "RECOMMENDED — DEMAND-WEIGHTED SPLIT",
             size=11.5, color=ORANGE_DARK, bold=True, char_spacing=1.2)
    _textbox(s2, left_x + 0.4, panel_y + 0.75, left_w - 0.8, 0.6, f"${recommended_cost:,.2f}",
             size=36, color=CHARCOAL, bold=True)
    _textbox(s2, left_x + 0.4, panel_y + 1.45, left_w - 0.8, 0.3, "ESTIMATED TOTAL LANDED COST",
             size=10, color=SLATE, char_spacing=1)
    # The recommended split's own "miss" is whatever share of demand sits in
    # regions it allocated zero units to (usually 0, since it spreads
    # proportionally across every region with nonzero demand -- but compute
    # it for real rather than assuming).
    active_regions = [r for r, u in region_split.items() if u > 0]
    recommended_missed_pct = sum(
        demand_pct.get(r, 0) for r in demand_pct if r not in active_regions
    ) * 100

    _textbox(s2, left_x + 0.4, panel_y + 2.05, left_w - 0.8, 1.8,
             f"Ships to all regions in proportion to where demand actually exists, "
             f"missing approximately {recommended_missed_pct:.0f}% of identified historical demand.",
             size=12.5, color=CHARCOAL, line_spacing=1.3)

    # Right panel: cheapest single-region
    _rect(s2, right_x, panel_y, right_w, panel_h, WHITE, line_color=CHARCOAL, shadow=True)
    _textbox(s2, right_x + 0.4, panel_y + 0.35, right_w - 0.8, 0.3,
             f"CHEAPEST OPTION — {cheapest_region.upper()} ONLY", size=11.5, color=SLATE, bold=True, char_spacing=1.2)
    _textbox(s2, right_x + 0.4, panel_y + 0.75, right_w - 0.8, 0.6, f"${cheapest_cost:,.2f}",
             size=36, color=CHARCOAL, bold=True)
    _textbox(s2, right_x + 0.4, panel_y + 1.45, right_w - 0.8, 0.3, "ESTIMATED TOTAL LANDED COST",
             size=10, color=SLATE, char_spacing=1)
    _textbox(s2, right_x + 0.4, panel_y + 2.05, right_w - 0.8, 1.8,
             f"Cheapest in freight + placement fees, but misses an estimated "
             f"{coverage_gap_pct:.0f} percentage points of this SKU's historical demand "
             f"by only stocking {cheapest_region}.",
             size=12.5, color=CHARCOAL, line_spacing=1.3)

    delta_label = "MORE" if cost_delta > 0 else ("LESS" if cost_delta < 0 else "THE SAME")
    delta_color = ORANGE_DARK if cost_delta > 0 else CHARCOAL
    _textbox(s2, 0.6, panel_y + panel_h + 0.25, 12.1, 0.5,
             f"DEMAND-WEIGHTED SPLIT COSTS ${abs(cost_delta):,.2f} {delta_label} THAN CHEAPEST-ONLY",
             size=14, color=delta_color, bold=True, char_spacing=0.5, align=PP_ALIGN.CENTER)

    _footer(s2, "Cost Tradeoff", "2", W, H)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
